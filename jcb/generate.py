#!/usr/bin/env python3
"""
JCB成果報告レポート生成ツール

BigQueryから出力したCSV（4種類）を元に、クライアントごとの報告スライドを
1つのPPTXファイルにまとめて生成する。

毎回の作業は独立したタスクフォルダ（tasks/YYYYMMDD/）に分離される。
CSVはtasks/YYYYMMDD/input/にコピーされ、PPTXはtasks/YYYYMMDD/output/に出力される。

使い方:
    # CSVの場所を指定して実行（タスクフォルダは自動作成）
    python generate.py 202602 --input-dir /path/to/csvs

    # CSVファイルを直接指定
    python generate.py 202602 file1.csv file2.csv file3.csv file4.csv
"""

import argparse
import copy
import datetime
import re
import shutil
import subprocess
import sys
import tempfile
import unicodedata
from pathlib import Path

# shared モジュールのインポート（report-builder/shared/ から読み込み）
sys.path.insert(0, str(Path(__file__).parent.parent))

from shared.csv_utils import load_and_group, detect_csv_files
from shared.chart_generator import generate_bar_chart, generate_double_donut_chart

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from config import (
    TEMPLATE_PATH, TASKS_DIR,
    TABLE_COLUMNS, TABLE_MAX_ROWS,
    CSV_PATTERNS, COLUMN_MAP, BRAND_DONUT_MAX,
)

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SHAPE_TAGS = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}

# タイトル表示用: 会社名から除去するサフィックス
_COMPANY_STRIP_RE = re.compile(r"（申込企業：[^）]+）")

# デフォルトのセルマージン (EMU): 左右各 91440 (0.1インチ)
_DEFAULT_CELL_MARGIN = 91440


def _estimate_text_width(text: str) -> float:
    """テキストの推定幅を返す（全角=1.0, 半角=0.5 単位）"""
    width = 0.0
    for c in text:
        if unicodedata.east_asian_width(c) in ("F", "W"):
            width += 1.0
        else:
            width += 0.5
    return width


def _calc_font_size_for_cell(
    text: str, col_width_emu: int, base_size_pt: float = 10.5, min_size_pt: float = 7.0
) :
    """セル幅に1行で収まるフォントサイズを返す。base_sizeで収まればNone。"""
    effective_width = col_width_emu - _DEFAULT_CELL_MARGIN * 2
    text_width = _estimate_text_width(text)
    if text_width == 0:
        return None

    # base_sizeで収まるか確認
    base_emu = base_size_pt * 12700
    if text_width * base_emu <= effective_width:
        return None  # 縮小不要

    # 0.5pt刻みで縮小して収まるサイズを探す
    size_pt = base_size_pt - 0.5
    while size_pt >= min_size_pt:
        font_emu = size_pt * 12700
        if text_width * font_emu <= effective_width:
            return size_pt
        size_pt -= 0.5

    return min_size_pt


def _set_cell_text(cell, text: str, font_size_pt=None):
    """テーブルセルのテキストを書式を維持したまま置換

    段落間隔も明示的に0に設定し、行高さの自動膨張を防ぐ。
    font_size_pt が指定された場合、そのサイズに変更する（改行回避用）。
    """
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = text
            if font_size_pt is not None:
                para.runs[0].font.size = Pt(font_size_pt)
            _ensure_compact_spacing(para)
            return
    # runがない場合（空セル）はそのまま
    if text:
        cell.text_frame.paragraphs[0].text = text
        _ensure_compact_spacing(cell.text_frame.paragraphs[0])


def _ensure_compact_spacing(para):
    """段落にspcBef=0, spcAft=0, lnSpc=100%を明示設定（行高さ膨張防止）"""
    from lxml import etree as _et
    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pPr = para._p.find(f"{{{_A}}}pPr")
    if pPr is None:
        pPr = _et.SubElement(para._p, f"{{{_A}}}pPr")
        para._p.insert(0, pPr)
    for tag, child_tag, attr, val in [
        ("spcBef", "spcPts", "val", "0"),
        ("spcAft", "spcPts", "val", "0"),
        ("lnSpc", "spcPct", "val", "100000"),
    ]:
        el = pPr.find(f"{{{_A}}}{tag}")
        if el is None:
            el = _et.SubElement(pPr, f"{{{_A}}}{tag}")
        # 既存の子を削除して再作成
        for child in list(el):
            el.remove(child)
        _et.SubElement(el, f"{{{_A}}}{child_tag}").set(attr, val)


def replace_text(shape, new_text: str):
    """シェイプのテキストを置換（書式を維持）"""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            return  # 最初のparagraphの最初のrunに入れたら終了


def replace_chart_label_runs(shape, period_month: str, value, unit: str = ""):
    """チャートラベルのrun単位テキスト置換（混在フォントサイズ維持）

    テンプレート構造:
      run0: "■ラベル名" (14pt)
      run1: "（" (12pt)
      run2: "yyyy" (12pt)
      run3: "/mm" (12pt)
      run4: "の〇〇" (12pt)
      run5: "：　人）" or "：　件）" (12pt) → "：{value}{unit}）" に置換
    """
    if not shape.has_text_frame:
        return
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    if len(runs) < 6:
        return

    # yyyy → 実際の年
    year, month = period_month.split("/") if "/" in period_month else ("yyyy", "mm")
    runs[2].text = year
    runs[3].text = f"/{month}"

    # 値を埋める（単位を保持）
    val_str = f"{int(value):,}" if isinstance(value, (int, float)) and value != "" else str(value)
    runs[5].text = f"：{val_str}{unit}）"


def replace_image_by_ref(slide, shape, image_path: Path):
    """シェイプ参照を使って画像を差し替え（位置・サイズ維持）"""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(image_path), left, top, width, height)


def fill_table_by_ref(shape, brand_df):
    """シェイプ参照を使ってテーブルにデータを書き込む"""
    if not shape.has_table:
        print(f"  警告: '{shape.name}' はテーブルではありません")
        return

    table = shape.table
    rows = list(table.rows)
    n_cols = len(list(rows[0].cells))

    # ブランド名列の幅を取得（フォント縮小判定用）
    tbl = table._tbl
    grid = tbl.find(f"{{{A_NS}}}tblGrid")
    col_widths = [int(col.get("w")) for col in grid.findall(f"{{{A_NS}}}gridCol")]
    brand_col_width = col_widths[0]

    # カラムごとの単位
    _COL_UNITS = {
        "total_price": "円",
        "total_count": "件",
        "unique_user_count": "人",
    }

    for i in range(TABLE_MAX_ROWS):
        row = rows[i + 1]
        if i < len(brand_df):
            data_row = brand_df.iloc[i]
            for j, col in enumerate(TABLE_COLUMNS):
                if j >= n_cols:
                    break
                val = data_row.get(col, "")
                if isinstance(val, float) and val != val:
                    val = ""
                elif isinstance(val, (int, float)) and col != "brand_name":
                    val = f"{int(val):,}"
                unit = _COL_UNITS.get(col, "")
                text = f"{val}{unit}" if val != "" else ""
                # ブランド名列のみ: 長い名前は折り返し防止のためフォント縮小
                font_size = None
                if j == 0 and text:
                    font_size = _calc_font_size_for_cell(text, brand_col_width)
                _set_cell_text(row.cells[j], text, font_size_pt=font_size)
            for j in range(len(TABLE_COLUMNS), n_cols):
                _set_cell_text(row.cells[j], "")
        else:
            for j in range(n_cols):
                _set_cell_text(row.cells[j], "")


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _copy_slide(src_prs, src_idx, dst_prs):
    """ソースPPTXのスライドをコピーしてdstに追加

    非画像シェイプはXML deepcopy、画像は一時ファイル経由で add_picture() を使う。
    これにより各スライドが一意なImagePartを持ち、ZIP内の名前衝突を回避する。
    """
    src_slide = src_prs.slides[src_idx]

    # ベースPPTXの1枚目と同じレイアウトを使用
    slide_layout = dst_prs.slides[0].slide_layout
    new_slide = dst_prs.slides.add_slide(slide_layout)

    # レイアウトから自動生成された全シェイプを削除（ソースからコピーするため）
    for el in list(new_slide.shapes._spTree):
        tag = etree.QName(el.tag).localname
        if tag in SHAPE_TAGS:
            new_slide.shapes._spTree.remove(el)

    dst_spTree = new_slide.shapes._spTree
    embed_attr = f"{{{R_NS}}}embed"

    for el in src_slide.shapes._spTree:
        tag = etree.QName(el.tag).localname
        if tag not in SHAPE_TAGS:
            continue

        if tag == "pic":
            # 画像シェイプ: blob→一時ファイル→add_picture（一意なImagePart生成）
            blip = el.find(f".//{{{A_NS}}}blip")
            if blip is not None:
                rId = blip.get(embed_attr)
                if rId and rId in src_slide.part.rels:
                    rel = src_slide.part.rels[rId]
                    blob = rel.target_part.blob
                    ct = rel.target_part.content_type
                    ext = ct.split("/")[-1].replace("jpeg", "jpg")

                    xfrm = el.find(f".//{{{A_NS}}}xfrm")
                    off = xfrm.find(f"{{{A_NS}}}off")
                    ext_el = xfrm.find(f"{{{A_NS}}}ext")
                    left = int(off.get("x", 0))
                    top = int(off.get("y", 0))
                    width = int(ext_el.get("cx", 0))
                    height = int(ext_el.get("cy", 0))

                    with tempfile.NamedTemporaryFile(
                        suffix=f".{ext}", delete=False
                    ) as f:
                        f.write(blob)
                        tmp_path = f.name
                    try:
                        new_slide.shapes.add_picture(
                            tmp_path, left, top, width, height
                        )
                    finally:
                        Path(tmp_path).unlink(missing_ok=True)
                    continue

        # 非画像シェイプ: XML deepcopy で書式を完全維持
        new_el = copy.deepcopy(el)
        dst_spTree.append(new_el)

    return new_slide


def _detect_shapes(slide):
    """テンプレートのシェイプをテキスト内容・タイプ・位置から自動検出する。

    PowerPointでテンプレートを編集するとシェイプのインデックスが変わるため、
    インデックスに依存せず、内容ベースで検出する。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    refs = {}
    text_shapes = []
    picture_shapes = []

    for shape in slide.shapes:
        if shape.has_table:
            refs["table"] = shape
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes.append(shape)
            continue
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "ご報告資料" in t:
                refs["title"] = shape
            elif "期間" in t and "yyyy" in t:
                refs["period"] = shape
            elif "初回登録ユーザー数" in t:
                refs["reg_users"] = shape
            elif "MAU" in t:
                refs["mau"] = shape
            elif "流通総額" in t or "総購入金額" in t:
                refs["distribution"] = shape
            elif "購入上位" in t or "TOP10" in t:
                refs["top_tickets_label"] = shape
            elif "購入数推移" in t:
                refs["purchase_chart_label"] = shape
            elif "ログインユーザー数推移" in t:
                refs["login_chart_label"] = shape
            elif "ブランド" in t:
                refs["brand_chart_label"] = shape
            elif t.strip() == "":
                # 空テキスト（右側のダミー）
                if shape.left > slide.shapes[0].left:
                    refs["empty_text"] = shape
            text_shapes.append(shape)

    # 画像シェイプをラベルとの位置関係で特定
    # 各ラベルの直下（y座標が近く、ラベルより下）にある画像を割り当て
    def _find_picture_below(label_key):
        if label_key not in refs:
            return None
        label = refs[label_key]
        label_top = label.top
        label_left = label.left
        best = None
        best_dist = float("inf")
        for pic in picture_shapes:
            # ラベルより下にある画像
            if pic.top > label_top:
                dist = pic.top - label_top
                if dist < best_dist:
                    best_dist = dist
                    best = pic
        return best

    refs["login_chart_image"] = _find_picture_below("login_chart_label")
    refs["purchase_chart_image"] = _find_picture_below("purchase_chart_label")

    # ブランド画像: ラベルの下にある画像（2つある場合は大きい方がimage1、小さい方がimage2）
    brand_label = refs.get("brand_chart_label")
    if brand_label:
        brand_pics = []
        for pic in picture_shapes:
            if pic.top > brand_label.top and pic not in (refs.get("login_chart_image"), refs.get("purchase_chart_image")):
                brand_pics.append(pic)
        # サイズ（幅）が大きい方をimage1、小さい方をimage2（重複削除対象）
        brand_pics.sort(key=lambda p: p.width, reverse=True)
        if len(brand_pics) >= 1:
            refs["brand_chart_image1"] = brand_pics[0]
        if len(brand_pics) >= 2:
            refs["brand_chart_image2"] = brand_pics[1]

    # 検出結果を表示
    detected = [k for k, v in refs.items() if v is not None]
    missing = [k for k in ["title", "period", "reg_users", "mau", "distribution",
               "table", "login_chart_image", "purchase_chart_image",
               "login_chart_label", "purchase_chart_label",
               "brand_chart_image1", "brand_chart_label"] if k not in detected]
    if missing:
        print(f"  警告: 未検出シェイプ: {', '.join(missing)}")

    return refs


def generate_single_report(
    project_id: str,
    project_name: str,
    summary_row,
    login_df,
    purchase_df,
    brand_df,
    tmp_dir: Path,
) -> str:
    """1クライアント分のPPTXを生成しパスを返す"""
    prs = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides[0]

    refs = _detect_shapes(slide)

    # --- テキスト置換 ---
    period_start = summary_row.get("period_start", "")
    period_end = summary_row.get("period_end", "")
    reg_users = summary_row.get("first_registration_users", "xx")
    mau = summary_row.get("mau", "xx")
    dist_total = summary_row.get("product_distribution_total", "xx")
    purchase_total = summary_row.get("total_purchase_amount", "xx")

    def fmt(v):
        try:
            return f"{int(v):,}"
        except (ValueError, TypeError):
            return str(v)

    display_name = _COMPANY_STRIP_RE.sub("", project_name)
    replace_text(refs["title"], f"{display_name} ご報告資料")
    replace_text(refs["period"], f"期間：{period_start}〜{period_end}")
    replace_text(refs["reg_users"], f"■初回登録ユーザー数：{fmt(reg_users)}人")
    replace_text(refs["mau"], f"■MAU(購入ユーザー数)：{fmt(mau)}人")

    dist_text = f"　商品代流通総額：{fmt(dist_total)}円 (総購入金額：{fmt(purchase_total)}円)"
    replace_text(refs["distribution"], dist_text)

    period_month = "/".join(period_start.split("/")[:2]) if period_start else "yyyy/mm"

    login_col = COLUMN_MAP["login"]["value"]
    purchase_col = COLUMN_MAP["purchase"]["value"]

    login_total = int(login_df[login_col].sum()) if login_df is not None and len(login_df) > 0 else ""
    purchase_total_count = int(purchase_df[purchase_col].sum()) if purchase_df is not None and len(purchase_df) > 0 else ""

    replace_chart_label_runs(refs["login_chart_label"], period_month, login_total, unit="人")
    replace_chart_label_runs(refs["purchase_chart_label"], period_month, purchase_total_count, unit="件")

    # --- テーブル ---
    # ブランドデータがある場合は書き込み、ない場合はテンプレのダミーをクリア
    fill_table_by_ref(refs["table"], brand_df if brand_df is not None and len(brand_df) > 0 else pd.DataFrame())

    # --- グラフ画像 ---
    # データがある場合: グラフ生成→画像差し替え
    # データがない場合: テンプレのダミー画像を削除
    login_date_col = COLUMN_MAP["login"]["date"]
    purchase_date_col = COLUMN_MAP["purchase"]["date"]

    def _remove_shape(ref_key):
        try:
            sp = refs[ref_key]._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    if login_df is not None and len(login_df) > 0:
        chart_path = tmp_dir / f"{project_id}_login.png"
        generate_bar_chart(
            dates=login_df[login_date_col].tolist(),
            values=login_df[login_col].tolist(),
            title="", ylabel="ユーザー数",
            output_path=chart_path,
            figsize=(5.77, 1.45), bar_width=5,
        )
        replace_image_by_ref(slide, refs["login_chart_image"], chart_path)
    else:
        _remove_shape("login_chart_image")

    if purchase_df is not None and len(purchase_df) > 0:
        chart_path = tmp_dir / f"{project_id}_purchase.png"
        generate_bar_chart(
            dates=purchase_df[purchase_date_col].tolist(),
            values=purchase_df[purchase_col].tolist(),
            title="", ylabel="購入数",
            output_path=chart_path,
            figsize=(5.77, 1.45), color="#34A853", bar_width=5,
        )
        replace_image_by_ref(slide, refs["purchase_chart_image"], chart_path)
    else:
        _remove_shape("purchase_chart_image")

    if brand_df is not None and len(brand_df) > 0:
        chart_path = tmp_dir / f"{project_id}_brand.png"
        donut_df = brand_df.head(BRAND_DONUT_MAX).copy()
        remaining = brand_df.iloc[BRAND_DONUT_MAX:]
        if len(remaining) > 0:
            other = pd.DataFrame([{
                "brand_name": "その他",
                "total_count": remaining["total_count"].sum(),
                "total_price": remaining["total_price"].sum(),
            }])
            donut_df = pd.concat([donut_df, other], ignore_index=True)

        generate_double_donut_chart(
            values1=donut_df["total_count"].tolist(),
            labels1=donut_df["brand_name"].tolist(),
            title1="発行総数",
            values2=donut_df["total_price"].tolist(),
            labels2=donut_df["brand_name"].tolist(),
            title2="販売総額",
            output_path=chart_path,
            figsize=(6.2, 1.55),
        )
        replace_image_by_ref(slide, refs["brand_chart_image1"], chart_path)
    else:
        _remove_shape("brand_chart_image1")

    # 重複ドーナツ画像を常に削除
    _remove_shape("brand_chart_image2")

    # 一時ファイルに保存
    out = tmp_dir / f"{hash(project_id)}.pptx"
    prs.save(str(out))
    return str(out)


def _create_task_dir() -> Path:
    """今日の日付でタスクフォルダを作成。同日に既存があればサフィックスを追加。"""
    today = datetime.date.today().strftime("%Y%m%d")
    task_dir = TASKS_DIR / today
    if task_dir.exists():
        n = 2
        while (TASKS_DIR / f"{today}_{n}").exists():
            n += 1
        task_dir = TASKS_DIR / f"{today}_{n}"
    task_dir.mkdir(parents=True)
    (task_dir / "input").mkdir()
    (task_dir / "output").mkdir()
    return task_dir


def main():
    parser = argparse.ArgumentParser(description="JCB成果報告レポート生成ツール")
    parser.add_argument("month", help="対象月 (YYYYMM形式)")
    parser.add_argument("csv_files", nargs="*", help="CSVファイルパス（直接指定）")
    parser.add_argument("--input-dir", type=Path, default=None, help="CSV入力ディレクトリ")

    args = parser.parse_args()

    if not TEMPLATE_PATH.exists():
        print(f"エラー: テンプレートが見つかりません: {TEMPLATE_PATH}")
        sys.exit(1)

    # タスクフォルダ作成
    task_dir = _create_task_dir()
    task_input = task_dir / "input"
    task_output = task_dir / "output"
    print(f"タスクフォルダ: {task_dir}")

    # CSVをタスクのinputにコピー
    if args.csv_files:
        # 直接指定されたCSVファイルをコピー
        for csv_path in args.csv_files:
            src = Path(csv_path)
            if not src.exists():
                print(f"エラー: ファイルが見つかりません: {src}")
                sys.exit(1)
            shutil.copy(str(src), str(task_input / src.name))
        print(f"CSVコピー完了: {len(args.csv_files)}ファイル → {task_input}")
    elif args.input_dir:
        # ディレクトリからCSVをコピー
        for csv_file in args.input_dir.glob("*.csv"):
            shutil.copy(str(csv_file), str(task_input / csv_file.name))
        print(f"CSVコピー完了: {args.input_dir} → {task_input}")
    else:
        print("エラー: CSVファイルまたは--input-dirを指定してください")
        sys.exit(1)

    # CSV自動検出（タスクのinputから）
    csv_files = detect_csv_files(task_input, CSV_PATTERNS)
    missing = [k for k, v in csv_files.items() if v is None]
    if missing:
        print(f"エラー: 以下のCSVが見つかりません: {', '.join(missing)}")
        print(f"  検索先: {task_input}")
        for k, patterns in CSV_PATTERNS.items():
            status = "OK" if csv_files[k] else "見つからない"
            print(f"    {k}: {patterns} → {status}")
        sys.exit(1)

    print("\n検出されたCSV:")
    for k, v in csv_files.items():
        print(f"  {k}: {v.name}")

    # --- CSVロード ---
    summary_name_col = COLUMN_MAP["summary"]["company_name"]
    brand_name_col = COLUMN_MAP["brand"]["company_name"]

    # サマリ: すでにクライアント単位で集約済み
    summary_df_all = pd.read_csv(csv_files["summary"])
    summary_by_client = {}
    for _, row in summary_df_all.iterrows():
        summary_by_client[row[summary_name_col]] = row

    # ログイン/購入: company_nameでグルーピング
    login_by_company = load_and_group(csv_files["login"], COLUMN_MAP["login"]["company_name"])
    purchase_by_company = load_and_group(csv_files["purchase"], COLUMN_MAP["purchase"]["company_name"])

    # 購入データの欠損週を0で埋める（ログインの全週を基準）
    login_date_col = COLUMN_MAP["login"]["date"]
    purchase_date_col = COLUMN_MAP["purchase"]["date"]
    purchase_val_col = COLUMN_MAP["purchase"]["value"]
    all_weeks = sorted(
        pd.read_csv(csv_files["login"])[login_date_col].unique()
    )
    for company, pdf in purchase_by_company.items():
        existing = set(pdf[purchase_date_col].values)
        missing = [w for w in all_weeks if w not in existing]
        if missing:
            fill = pd.DataFrame({purchase_date_col: missing})
            for col in pdf.columns:
                if col != purchase_date_col:
                    fill[col] = 0 if col == purchase_val_col else pdf.iloc[0].get(col, "")
            pdf = pd.concat([pdf, fill], ignore_index=True)
            pdf = pdf.sort_values(purchase_date_col).reset_index(drop=True)
            purchase_by_company[company] = pdf

    # ブランド: クライアント別にグルーピング
    brand_df_all = pd.read_csv(csv_files["brand"])
    brand_by_company = {}
    for company, bdf in brand_df_all.groupby(brand_name_col):
        brand_by_company[company] = bdf.sort_values(
            "total_price", ascending=False
        ).reset_index(drop=True)

    client_names = list(summary_by_client.keys())
    print(f"\nクライアント数: {len(client_names)}")

    # 1. 各クライアントの個別PPTXを一時生成
    # 2. 全スライドを1ファイルに統合
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        temp_pptx_paths = []

        for client_name in client_names:
            summary_row = summary_by_client[client_name]
            login_df = login_by_company.get(client_name)
            purchase_df = purchase_by_company.get(client_name)
            client_brand_df = brand_by_company.get(client_name)

            print(f"\n生成中: {client_name}")
            if login_df is None:
                print("  注意: ログインデータなし")
            if purchase_df is None:
                print("  注意: 購入データなし")
            if client_brand_df is None:
                print("  注意: ブランドデータなし")

            pptx_path = generate_single_report(
                project_id=client_name,
                project_name=client_name,
                summary_row=summary_row,
                login_df=login_df,
                purchase_df=purchase_df,
                brand_df=client_brand_df,
                tmp_dir=tmp_path,
            )
            temp_pptx_paths.append(pptx_path)
            print(f"  → スライド {len(temp_pptx_paths)}")

        # 統合: 最初のPPTXをベースに残りのスライドを追加
        print("\nスライドを統合中...")
        base_prs = Presentation(temp_pptx_paths[0])

        for i, path in enumerate(temp_pptx_paths[1:], 2):
            src_prs = Presentation(path)
            _copy_slide(src_prs, 0, base_prs)

        output_file = task_output / f"JCB報告資料_{args.month}.pptx"
        base_prs.save(str(output_file))

    print(f"\n完了! 出力先: {output_file}")
    print(f"  タスクフォルダ: {task_dir}")
    print(f"  スライド数: {len(client_names)}")

    # 生成したPPTXを自動で開く
    subprocess.Popen(["open", str(output_file)])


if __name__ == "__main__":
    main()
