#!/usr/bin/env python3
"""JCBレポートのバリデーション

生成結果のPPTXをテンプレートと比較し、テーブル構造が崩れていないか検証する。

使い方:
    python validate.py tasks/20260220/output/JCB報告資料_202602.pptx
"""

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from config import TEMPLATE_PATH, SHAPES

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _get_table_metrics(slide):
    """スライドからテーブルのメトリクスを取得"""
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            tbl = table._tbl
            rows = list(table.rows)

            row_heights = []
            for row in rows:
                h = row._tr.get("h", "0")
                row_heights.append(int(h))

            grid = tbl.find(f"{{{A_NS}}}tblGrid")
            col_widths = [int(c.get("w")) for c in grid.findall(f"{{{A_NS}}}gridCol")]

            # 各セルの段落数を確認
            para_counts = []
            for row in rows:
                row_paras = []
                for cell in row.cells:
                    row_paras.append(len(cell.text_frame.paragraphs))
                para_counts.append(row_paras)

            return {
                "top": shape.top,
                "left": shape.left,
                "width": shape.width,
                "height": shape.height,
                "n_rows": len(rows),
                "n_cols": len(col_widths),
                "row_heights": row_heights,
                "col_widths": col_widths,
                "para_counts": para_counts,
            }
    return None


def validate(pptx_path: str) -> list[dict]:
    """生成結果を検証し、問題リストを返す"""
    issues = []

    # テンプレートのメトリクス
    tmpl_prs = Presentation(str(TEMPLATE_PATH))
    tmpl_metrics = _get_table_metrics(tmpl_prs.slides[0])

    # 生成結果
    gen_prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(gen_prs.slides):
        slide_num = slide_idx + 1
        gen_metrics = _get_table_metrics(slide)

        if gen_metrics is None:
            issues.append({"slide": slide_num, "severity": "error", "message": "テーブルが見つかりません"})
            continue

        # 行数チェック
        if gen_metrics["n_rows"] != tmpl_metrics["n_rows"]:
            issues.append({
                "slide": slide_num, "severity": "error",
                "message": f"行数不一致: テンプレ={tmpl_metrics['n_rows']}, 生成={gen_metrics['n_rows']}",
            })

        # 列数チェック
        if gen_metrics["n_cols"] != tmpl_metrics["n_cols"]:
            issues.append({
                "slide": slide_num, "severity": "error",
                "message": f"列数不一致: テンプレ={tmpl_metrics['n_cols']}, 生成={gen_metrics['n_cols']}",
            })

        # 行高さチェック
        for i, (th, gh) in enumerate(zip(tmpl_metrics["row_heights"], gen_metrics["row_heights"])):
            if th != gh:
                issues.append({
                    "slide": slide_num, "severity": "warning",
                    "message": f"行{i}の高さ不一致: テンプレ={th}, 生成={gh}",
                })

        # 列幅チェック
        for i, (tw, gw) in enumerate(zip(tmpl_metrics["col_widths"], gen_metrics["col_widths"])):
            if tw != gw:
                issues.append({
                    "slide": slide_num, "severity": "error",
                    "message": f"列{i}の幅不一致: テンプレ={tw}, 生成={gw}",
                })

        # テーブル位置チェック
        if gen_metrics["top"] != tmpl_metrics["top"] or gen_metrics["left"] != tmpl_metrics["left"]:
            issues.append({
                "slide": slide_num, "severity": "warning",
                "message": f"テーブル位置ずれ: テンプレ=({tmpl_metrics['left']},{tmpl_metrics['top']}), 生成=({gen_metrics['left']},{gen_metrics['top']})",
            })

        # テーブルサイズチェック
        if gen_metrics["width"] != tmpl_metrics["width"] or gen_metrics["height"] != tmpl_metrics["height"]:
            issues.append({
                "slide": slide_num, "severity": "error",
                "message": f"テーブルサイズ不一致: テンプレ=({tmpl_metrics['width']}x{tmpl_metrics['height']}), 生成=({gen_metrics['width']}x{gen_metrics['height']})",
            })

        # 段落数チェック（各セルが1段落であること）
        for i, row_paras in enumerate(gen_metrics["para_counts"]):
            for j, n_paras in enumerate(row_paras):
                if n_paras > 1:
                    issues.append({
                        "slide": slide_num, "severity": "error",
                        "message": f"行{i}列{j}: 段落数が{n_paras}（1であるべき）",
                    })

    return issues


def main():
    if len(sys.argv) < 2:
        print("使い方: python validate.py <pptx_path>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not Path(pptx_path).exists():
        print(f"エラー: ファイルが見つかりません: {pptx_path}")
        sys.exit(1)

    issues = validate(pptx_path)

    if not issues:
        print("OK: テンプレートと構造が一致しています")
        sys.exit(0)

    errors = [i for i in issues if i["severity"] == "error"]
    warnings = [i for i in issues if i["severity"] == "warning"]

    if errors:
        print(f"ERROR: {len(errors)}件の問題が見つかりました")
        for issue in errors:
            print(f"  スライド{issue['slide']}: {issue['message']}")

    if warnings:
        print(f"WARNING: {len(warnings)}件の警告があります")
        for issue in warnings:
            print(f"  スライド{issue['slide']}: {issue['message']}")

    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
